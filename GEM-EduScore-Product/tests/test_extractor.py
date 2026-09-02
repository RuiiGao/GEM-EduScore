"""Document-ingestion tests for supported Education material formats."""

from __future__ import annotations

from io import BytesIO
import unittest

from modules.extractor import extract_document_text


class ExtractorTests(unittest.TestCase):
    def test_docx_paragraphs_and_tables_are_extracted(self) -> None:
        from docx import Document

        document = Document()
        document.add_heading("Education Workshop", level=1)
        document.add_paragraph("Students built a cell model and discussed bioethics.")
        table = document.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "Audience"
        table.cell(0, 1).text = "Secondary students"
        stream = BytesIO()
        document.save(stream)

        text, info = extract_document_text("portfolio.docx", stream.getvalue())
        self.assertIn("Education Workshop", text)
        self.assertIn("Audience | Secondary students", text)
        self.assertEqual(info["extension"], "DOCX")

    def test_text_pdf_is_extracted_with_page_markers(self) -> None:
        from reportlab.pdfgen import canvas

        stream = BytesIO()
        pdf = canvas.Canvas(stream)
        pdf.drawString(72, 760, "Education outcome survey and participant feedback")
        pdf.showPage()
        pdf.save()

        text, info = extract_document_text("evidence.pdf", stream.getvalue())
        self.assertIn("[PDF Page 1]", text)
        self.assertIn("participant feedback", text)
        self.assertEqual(info["pages_or_slides"], 1)

    def test_pptx_slide_text_and_table_are_extracted(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Interactive Education"
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(7), Inches(1)).table
        table.cell(0, 0).text = "Method"
        table.cell(0, 1).text = "Pre/post questionnaire"
        stream = BytesIO()
        presentation.save(stream)

        text, info = extract_document_text("slides.pptx", stream.getvalue())
        self.assertIn("[Slide 1]", text)
        self.assertIn("Pre/post questionnaire", text)
        self.assertEqual(info["pages_or_slides"], 1)

    def test_html_ignores_script_content(self) -> None:
        raw = b"<h1>Education Plan</h1><script>ignore_me()</script><p>Workshop evidence</p>"
        text, _ = extract_document_text("page.html", raw)
        self.assertIn("Workshop evidence", text)
        self.assertNotIn("ignore_me", text)

    def test_legacy_doc_has_actionable_message(self) -> None:
        with self.assertRaisesRegex(ValueError, "另存为 .docx"):
            extract_document_text("legacy.doc", b"old word content")


if __name__ == "__main__":
    unittest.main()
