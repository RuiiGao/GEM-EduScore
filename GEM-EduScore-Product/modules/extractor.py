"""Local document ingestion and demo evidence helpers.

This module deliberately contains no LLM or external API integration. The first
prototype validates local files and exposes a stable evidence shape for the UI.
"""

from __future__ import annotations

from dataclasses import asdict, dataclass
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path
from typing import Any


SUPPORTED_EXTENSIONS = {".md", ".txt", ".docx", ".pdf", ".pptx", ".html", ".htm", ".csv"}
MAX_FILE_SIZE_BYTES = 20 * 1024 * 1024


@dataclass(frozen=True)
class DocumentInfo:
    name: str
    extension: str
    characters: int
    words: int
    lines: int
    preview: str
    pages_or_slides: int | None = None


def extract_document_text(file_name: str, raw_bytes: bytes) -> tuple[str, dict[str, Any]]:
    """Extract readable text from supported education-material formats."""
    extension = Path(file_name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        if extension == ".doc":
            raise ValueError("暂不支持旧版 .doc 文件。请在 Word 中另存为 .docx 后上传。")
        raise ValueError(
            "不支持该文件格式。当前支持 MD、TXT、DOCX、PDF、PPTX、HTML 和 CSV。"
        )
    if not raw_bytes:
        raise ValueError("上传的文件为空，请选择包含内容的材料。")
    if len(raw_bytes) > MAX_FILE_SIZE_BYTES:
        raise ValueError("文件超过 20 MB。请压缩图片、拆分文件或上传精简版本。")

    pages_or_slides: int | None = None
    if extension in {".md", ".txt", ".csv"}:
        text = _decode_text(raw_bytes)
    elif extension in {".html", ".htm"}:
        text = _extract_html(raw_bytes)
    elif extension == ".docx":
        text = _extract_docx(raw_bytes)
    elif extension == ".pdf":
        text, pages_or_slides = _extract_pdf(raw_bytes)
    else:
        text, pages_or_slides = _extract_pptx(raw_bytes)

    text = _clean_text(text)
    if not text:
        raise ValueError("没有从文件中读取到有效文本。请确认文件不是空白或纯扫描图片。")

    info = DocumentInfo(
        name=file_name,
        extension=extension.removeprefix(".").upper(),
        characters=len(text),
        words=len(text.split()),
        lines=len(text.splitlines()),
        preview=text[:900],
        pages_or_slides=pages_or_slides,
    )
    return text, asdict(info)


def get_demo_evidence_profile() -> dict[str, Any]:
    """Return the hand-curated JLU-CP profile used to exercise the UI."""
    return {
        "practice_name": "JLU-CP 2025 Education Portfolio",
        "summary": (
            "面向儿童、学生、老年群体与公众的多层次合成生物学教育实践，"
            "通过讲座、互动游戏、模型制作和创造性任务连接科学知识与实践。"
        ),
        "goals": [
            "激发儿童对生命科学的兴趣与好奇心",
            "构建覆盖不同生命阶段的合成生物学终身学习框架",
        ],
        "audiences": ["儿童与家长", "学生", "老年群体", "公众"],
        "design": [
            "健康讲座与知识输入",
            "互动问答和小游戏",
            "细胞模型制作",
            "细胞设计创意竞赛",
            "Synthetic Eco-Town 系统设计任务",
        ],
        "strong_evidence": [
            {
                "title": "受众与目标明确",
                "detail": "材料明确说明儿童与家长等目标群体，并给出激发生命科学兴趣的教育目标。",
                "source": "Children and parents / stimulate children's interest and curiosity in life sciences",
            },
            {
                "title": "形成完整学习链",
                "detail": "活动从知识输入进入互动、动手实践，再推进到创造性设计任务。",
                "source": "Health Lecture → Q&A and Games → Cell Model Making → Creative Design",
            },
            {
                "title": "存在参与者产出",
                "detail": "参与者制作细胞模型，并完成微生物解决方案等设计成果。",
                "source": "Cell models and microbial solution designs",
            },
        ],
        "missing_evidence": [
            "缺少与教育目标对应的前测、后测或学习变化数据",
            "缺少反馈如何推动下一轮活动修改的闭环记录",
            "缺少可供其他团队直接复用的完整教学包与实施说明",
            "缺少长期参与、跟踪和社区维护证据",
            "缺少生物伦理、技术风险与社会影响讨论记录",
        ],
        "coverage_note": "高质量活动也可能因文档记录不足而具有较低证据覆盖率。",
    }


def _decode_text(raw_bytes: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("无法识别文件编码，请将文件保存为 UTF-8 后重试。")


def _extract_docx(raw_bytes: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise ValueError("缺少 DOCX 解析组件。请重新运行双击启动器以安装依赖。") from exc

    try:
        document = Document(BytesIO(raw_bytes))
    except Exception as exc:
        raise ValueError("无法读取该 DOCX。文件可能已损坏、加密或不是有效的 Word 文档。") from exc

    parts: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    for table_index, table in enumerate(document.tables, 1):
        parts.append(f"\n[Table {table_index}]")
        for row in table.rows:
            cells = [_clean_inline(cell.text) for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pdf(raw_bytes: bytes) -> tuple[str, int]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValueError("缺少 PDF 解析组件。请重新运行双击启动器以安装依赖。") from exc

    try:
        reader = PdfReader(BytesIO(raw_bytes))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception as exc:
                raise ValueError("该 PDF 已加密，请先解除密码保护后上传。") from exc
        page_texts = []
        for page_number, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                page_texts.append(f"[PDF Page {page_number}]\n{text}")
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("无法读取该 PDF。文件可能已损坏或使用了不支持的编码。") from exc

    if not page_texts:
        raise ValueError("该 PDF 没有可提取的文本层；如果是扫描件，请先进行 OCR 后再上传。")
    return "\n\n".join(page_texts), len(reader.pages)


def _extract_pptx(raw_bytes: bytes) -> tuple[str, int]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("缺少 PPTX 解析组件。请重新运行双击启动器以安装依赖。") from exc

    try:
        presentation = Presentation(BytesIO(raw_bytes))
    except Exception as exc:
        raise ValueError("无法读取该 PPTX。文件可能已损坏或不是有效的 PowerPoint 文档。") from exc

    slides: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        content: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                content.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [_clean_inline(cell.text) for cell in row.cells]
                    if any(cells):
                        content.append(" | ".join(cells))
        if content:
            slides.append(f"[Slide {slide_number}]\n" + "\n".join(content))
    if not slides:
        raise ValueError("该 PPTX 中没有可提取的文字；图片中的文字需要先进行 OCR。")
    return "\n\n".join(slides), len(presentation.slides)


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript"}:
            self._ignored_depth += 1
        elif tag in {"p", "br", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"} and self._ignored_depth:
            self._ignored_depth -= 1

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth and data.strip():
            self.parts.append(data)


def _extract_html(raw_bytes: bytes) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(_decode_text(raw_bytes))
    return " ".join(parser.parts)


def _clean_inline(text: str) -> str:
    return " ".join(text.split())


def _clean_text(text: str) -> str:
    lines = [_clean_inline(line) for line in text.replace("\x00", "").splitlines()]
    cleaned: list[str] = []
    previous_blank = False
    for line in lines:
        if line:
            cleaned.append(line)
            previous_blank = False
        elif cleaned and not previous_blank:
            cleaned.append("")
            previous_blank = True
    return "\n".join(cleaned).strip()
